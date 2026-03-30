from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_PPTX = Path(r"C:\Users\bhalo\Downloads\_EXT_ Solution Challenge 2026 _ Prototype PPT Template.pptx")
OUTPUT_PPTX = ROOT / "submission_assets" / "JanHelp_Solution_Challenge_2026_15_Slides.pptx"
PREVIEW_DIR = ROOT / "submission_assets" / "preview"
LOGO_PATH = ROOT / "smartcity_application" / "assets" / "images" / "logo.png"


TEAM_NAME = "JanHelp"
TEAM_LEADER = "Kartik Bhalodiya"
TAGLINE = "AI-enabled civic complaint platform"

IMAGE_SLIDE_MAP = {
    2: "slide-01.png",
    3: "slide-02.png",
    4: "slide-03.png",
    5: "slide-04.png",
    6: "slide-05.png",
    8: "slide-06.png",
    9: "slide-07.png",
    11: "slide-08.png",
    12: "slide-09.png",
    13: "slide-10.png",
}


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.replace("#", "")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


BLUE = rgb("#2563EB")
BLUE_DARK = rgb("#1D4ED8")
GREEN = rgb("#16A34A")
AMBER = rgb("#F59E0B")
PURPLE = rgb("#7C3AED")
SLATE = rgb("#334155")
SLATE_LIGHT = rgb("#64748B")
SLATE_SOFT = rgb("#E2E8F0")
SLATE_BG = rgb("#F8FAFC")
SKY_SOFT = rgb("#E0F2FE")
GREEN_SOFT = rgb("#DCFCE7")
AMBER_SOFT = rgb("#FEF3C7")
PURPLE_SOFT = rgb("#EDE9FE")
TEAL_SOFT = rgb("#CCFBF1")
WHITE = rgb("#FFFFFF")


def add_full_slide_image(prs: Presentation, slide_index: int, image_name: str) -> None:
    image_path = PREVIEW_DIR / image_name
    if not image_path.exists():
        raise FileNotFoundError(f"Missing slide image: {image_path}")
    slide = prs.slides[slide_index - 1]
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)


def add_rounded_box(slide, x, y, w, h, fill_color, line_color=SLATE_SOFT, radius_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    return shape


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    font_size=18,
    color=SLATE,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Aptos",
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def add_bullet_list(slide, x, y, w, bullets, color=SLATE, bullet_color=BLUE, font_size=16):
    top = y
    for bullet in bullets:
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, top + Inches(0.08), Inches(0.08), Inches(0.08))
        circle.fill.solid()
        circle.fill.fore_color.rgb = bullet_color
        circle.line.fill.background()
        add_textbox(slide, x + Inches(0.15), top, w - Inches(0.15), Inches(0.45), bullet, font_size=font_size, color=color)
        top += Inches(0.5)


def build_wireframe_slide(slide) -> None:
    header = add_rounded_box(slide, Inches(0.55), Inches(1.55), Inches(3.0), Inches(0.55), SKY_SOFT)
    header.line.fill.background()
    add_textbox(slide, Inches(0.8), Inches(1.64), Inches(2.5), Inches(0.3), "Low-fidelity mobile wireframes", font_size=18, color=BLUE_DARK, bold=True)

    phone_specs = [
        (Inches(0.95), Inches(2.18), "Home"),
        (Inches(3.75), Inches(2.18), "AI Assistant"),
        (Inches(6.55), Inches(2.18), "Track Status"),
    ]

    for x, y, label in phone_specs:
        body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(1.85), Inches(2.85))
        body.fill.solid()
        body.fill.fore_color.rgb = WHITE
        body.line.color.rgb = SLATE
        body.line.width = Pt(3)

        notch = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(0.55), y + Inches(0.05), Inches(0.75), Inches(0.08))
        notch.fill.solid()
        notch.fill.fore_color.rgb = SLATE_SOFT
        notch.line.fill.background()

        add_textbox(slide, x + Inches(0.45), y - Inches(0.35), Inches(1.0), Inches(0.2), label, font_size=14, color=SLATE, bold=True, align=PP_ALIGN.CENTER)

        # Header bar
        top_bar = add_rounded_box(slide, x + Inches(0.15), y + Inches(0.18), Inches(1.55), Inches(0.42), WHITE, line_color=SLATE_SOFT)
        add_textbox(slide, x + Inches(0.22), y + Inches(0.22), Inches(1.2), Inches(0.2), "JanHelp", font_size=12, color=SLATE, bold=True)

        if label == "Home":
            for i, name in enumerate(["Live Stats", "Departments", "Quick Report"]):
                card = add_rounded_box(slide, x + Inches(0.22), y + Inches(0.78 + i * 0.55), Inches(1.42), Inches(0.4), WHITE)
                add_textbox(slide, x + Inches(0.32), y + Inches(0.88 + i * 0.55), Inches(0.9), Inches(0.15), name, font_size=10, color=SLATE, bold=True)
        elif label == "AI Assistant":
            bubble_colors = [SKY_SOFT, GREEN_SOFT, SKY_SOFT, PURPLE_SOFT]
            for i, fill in enumerate(bubble_colors):
                bubble = add_rounded_box(slide, x + Inches(0.3 + (0.15 if i % 2 else 0)), y + Inches(0.88 + i * 0.42), Inches(1.0), Inches(0.28), fill, line_color=fill)
                bubble.line.fill.background()
            add_rounded_box(slide, x + Inches(0.22), y + Inches(2.38), Inches(1.42), Inches(0.26), WHITE)
            add_textbox(slide, x + Inches(0.3), y + Inches(2.43), Inches(1.18), Inches(0.12), "Describe issue in any language", font_size=8, color=SLATE_LIGHT)
        else:
            labels = ["SC482913", "Assigned to Road Dept", "In progress", "Resolved + reopen"]
            dots = [BLUE, rgb("#0EA5E9"), AMBER, GREEN]
            for i, (text, dot) in enumerate(zip(labels, dots)):
                dot_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.22), y + Inches(0.95 + i * 0.5), Inches(0.08), Inches(0.08))
                dot_shape.fill.solid()
                dot_shape.fill.fore_color.rgb = dot
                dot_shape.line.fill.background()
                add_textbox(slide, x + Inches(0.42), y + Inches(0.92 + i * 0.5), Inches(1.05), Inches(0.15), text, font_size=9.5, color=SLATE, bold=True)

    footer = add_rounded_box(slide, Inches(1.6), Inches(5.08), Inches(6.8), Inches(0.34), WHITE)
    add_textbox(
        slide,
        Inches(1.78),
        Inches(5.12),
        Inches(6.4),
        Inches(0.2),
        "These mock screens represent the implemented citizen dashboard, AI assistant, complaint intake, and tracking flow.",
        font_size=11,
        color=SLATE,
        align=PP_ALIGN.CENTER,
    )


def build_cost_slide(slide) -> None:
    note = add_rounded_box(slide, Inches(0.55), Inches(1.55), Inches(3.3), Inches(0.55), AMBER_SOFT, line_color=AMBER_SOFT)
    add_textbox(slide, Inches(0.75), Inches(1.64), Inches(2.9), Inches(0.25), "Illustrative pilot estimate for a small-city rollout", font_size=16, color=AMBER, bold=True)

    table = add_rounded_box(slide, Inches(0.8), Inches(2.15), Inches(8.3), Inches(2.6), WHITE)
    rows = [
        ("Cloud hosting and deployment", "Rs 0 - 2,000 / month"),
        ("Database and storage", "Rs 0 - 3,500 / month"),
        ("Cloudinary media handling", "Rs 0 - 2,500 / month"),
        ("Gemini and AI usage", "Rs 1,000 - 4,000 / month"),
        ("Email, alerts, and maintenance buffer", "Rs 1,500 - 3,000 / month"),
    ]

    add_textbox(slide, Inches(1.05), Inches(2.35), Inches(4.2), Inches(0.25), "Cost Head", font_size=15, color=BLUE_DARK, bold=True)
    add_textbox(slide, Inches(6.25), Inches(2.35), Inches(2.3), Inches(0.25), "Estimated Cost", font_size=15, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)

    y = Inches(2.7)
    for left, right in rows:
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.0), y - Inches(0.04), Inches(7.5), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = SLATE_SOFT
        line.line.fill.background()
        add_textbox(slide, Inches(1.05), y, Inches(4.7), Inches(0.28), left, font_size=13.5, color=SLATE)
        add_textbox(slide, Inches(6.0), y, Inches(2.4), Inches(0.28), right, font_size=13.5, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
        y += Inches(0.38)

    total = add_rounded_box(slide, Inches(1.0), Inches(4.6), Inches(7.5), Inches(0.46), SKY_SOFT, line_color=SKY_SOFT)
    add_textbox(slide, Inches(1.2), Inches(4.68), Inches(4.5), Inches(0.2), "Estimated pilot total", font_size=15, color=BLUE_DARK, bold=True)
    add_textbox(slide, Inches(6.0), Inches(4.66), Inches(2.1), Inches(0.2), "Approx Rs 4,500 - 15,000 / month", font_size=15, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)

    footer = add_rounded_box(slide, Inches(1.4), Inches(5.15), Inches(6.7), Inches(0.32), WHITE)
    add_textbox(slide, Inches(1.55), Inches(5.2), Inches(6.35), Inches(0.16), "The MVP can begin on free or low-cost tiers and scale costs only when complaint volume grows.", font_size=11.5, color=SLATE, align=PP_ALIGN.CENTER)


def build_impact_slide(slide) -> None:
    add_rounded_box(slide, Inches(0.6), Inches(0.7), Inches(3.2), Inches(0.58), SKY_SOFT, line_color=SKY_SOFT)
    add_textbox(slide, Inches(0.82), Inches(0.81), Inches(2.8), Inches(0.22), "Pilot Readiness and Expected Impact", font_size=19, color=BLUE_DARK, bold=True)

    cards = [
        (Inches(0.8), Inches(1.55), Inches(2.65), Inches(1.2), SKY_SOFT, BLUE_DARK, "24/7 access", "Citizens can report issues anytime through multilingual guided intake."),
        (Inches(3.7), Inches(1.55), Inches(2.65), Inches(1.2), GREEN_SOFT, GREEN, "Cleaner data", "AI proof checks and duplicate suppression improve complaint quality."),
        (Inches(6.6), Inches(1.55), Inches(2.65), Inches(1.2), AMBER_SOFT, AMBER, "Faster routing", "Nearest-department assignment reduces manual triage and delay."),
    ]
    for x, y, w, h, fill, accent, title, body in cards:
        add_rounded_box(slide, x, y, w, h, fill, line_color=fill)
        add_textbox(slide, x + Inches(0.18), y + Inches(0.16), w - Inches(0.36), Inches(0.22), title, font_size=16, color=accent, bold=True)
        add_textbox(slide, x + Inches(0.18), y + Inches(0.46), w - Inches(0.36), Inches(0.55), body, font_size=12.5, color=SLATE)

    left = add_rounded_box(slide, Inches(0.95), Inches(3.05), Inches(3.85), Inches(1.8), WHITE)
    add_textbox(slide, Inches(1.2), Inches(3.2), Inches(3.2), Inches(0.25), "Why JanHelp is deployment-ready now", font_size=17, color=BLUE_DARK, bold=True)
    add_bullet_list(
        slide,
        Inches(1.15),
        Inches(3.55),
        Inches(3.3),
        [
            "Cloud deployment already defined in the backend stack",
            "Google Gemini usage already built into proof verification",
            "Citizen, department, and admin workflows already modeled",
        ],
        font_size=12.5,
    )

    right = add_rounded_box(slide, Inches(5.1), Inches(3.05), Inches(3.3), Inches(1.8), WHITE)
    add_textbox(slide, Inches(5.35), Inches(3.2), Inches(2.8), Inches(0.25), "Suggested pilot KPIs", font_size=17, color=PURPLE, bold=True)
    add_bullet_list(
        slide,
        Inches(5.3),
        Inches(3.55),
        Inches(2.8),
        [
            "Complaint submission time reduced",
            "Duplicate complaints reduced",
            "Routing accuracy improved",
            "Citizen satisfaction improved",
        ],
        font_size=12.5,
        bullet_color=PURPLE,
    )


def build_cover_slide(slide) -> None:
    panel = add_rounded_box(slide, Inches(0.85), Inches(2.75), Inches(3.6), Inches(1.6), WHITE)
    panel.fill.transparency = 0.08
    add_textbox(slide, Inches(1.1), Inches(3.0), Inches(3.0), Inches(0.45), "JanHelp", font_size=28, color=BLUE_DARK, bold=True)
    add_textbox(slide, Inches(1.12), Inches(3.45), Inches(3.2), Inches(0.3), TAGLINE, font_size=15, color=SLATE, bold=False)
    add_textbox(slide, Inches(1.12), Inches(3.78), Inches(3.1), Inches(0.25), f"Team: {TEAM_NAME}  |  Lead: {TEAM_LEADER}", font_size=12.5, color=SLATE_LIGHT)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(6.9), Inches(2.9), height=Inches(1.2))


def main() -> None:
    if not TEMPLATE_PPTX.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PPTX}")

    prs = Presentation(str(TEMPLATE_PPTX))

    for slide_index, image_name in IMAGE_SLIDE_MAP.items():
        add_full_slide_image(prs, slide_index, image_name)

    build_wireframe_slide(prs.slides[6])
    build_cost_slide(prs.slides[9])
    build_impact_slide(prs.slides[13])
    build_cover_slide(prs.slides[14])

    prs.save(str(OUTPUT_PPTX))
    print(f"Created: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
